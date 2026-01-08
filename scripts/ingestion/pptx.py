from hashlib import sha1
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from scripts.ingestion.models import AbstractIngestor, UnsupportedFileError
import logging
from scripts.utils.image_utils import (
    get_project_image_dir,
    infer_project_root,
    record_image_metadata,
)
from scripts.utils.logger import LoggerManager

# Initialize logger with proper file output
logger = LoggerManager.get_logger("pptx_ingestor")


class PptxIngestor(AbstractIngestor):
    """
    Ingestor for PPTX files.
    """

    def ingest(self, filepath: str) -> list[tuple[str, dict]]:
        if not filepath.endswith(".pptx"):
            raise UnsupportedFileError("File is not a .pptx file.")

        extracted_data = []
        seen_hashes = set()  # Track unique images

        try:
            prs = Presentation(filepath)
            file_path = Path(filepath)

            # 🔍 Infer project root and name
            project_root = infer_project_root(file_path)
            project_name = project_root.name
            image_dir = get_project_image_dir(project_name)

            logger.info(f"[PPTX] Ingesting file: {file_path}")
            logger.info(f"[PPTX] Image output folder: {image_dir}")

            file_stem = file_path.stem

            for i, slide in enumerate(prs.slides):
                slide_number = i + 1
                text_on_slide = []
                image_counter = 0
                image_paths = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_on_slide.append(shape.text.strip())

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        img_bytes = shape.image.blob
                        img_hash = sha1(img_bytes).hexdigest()
                        if img_hash in seen_hashes:
                            continue  # Skip duplicate image
                        seen_hashes.add(img_hash)

                        image_counter += 1
                        img_name = (
                            f"{file_stem}_slide{slide_number}_img{image_counter}.png"
                        )
                        rel_dir = Path("input") / "cache" / "images"
                        out_path = project_root / rel_dir / img_name
                        out_path.parent.mkdir(parents=True, exist_ok=True)

                        try:
                            with open(out_path, "wb") as f:
                                f.write(img_bytes)
                            logger.info(f"[PPTX] Saved image to: {out_path}")
                        except Exception as img_err:
                            logger.warning(
                                f"[PPTX] Failed to save image on slide "
                                f"{slide_number}: {img_err}"
                            )
                            continue

                        try:
                            out_path = out_path.resolve()
                            rel_to_input = out_path.relative_to(project_root / "input")
                            img_rel = str(rel_to_input)
                            image_paths.append(img_rel)
                        except Exception as e:
                            logger.warning(
                                f"[PPTX] Failed to compute relative image path: {e}"
                            )
                            continue

                # Emit main slide chunk (text + images)
                if text_on_slide or image_paths:
                    slide_meta = {
                        "slide_number": slide_number,
                        "type": "slide_content",
                        "doc_type": "pptx",
                    }
                    if image_paths:
                        slide_meta["image_paths"] = image_paths

                    text_content = (
                        "\n".join(text_on_slide).strip() or "[Image-only slide]"
                    )
                    extracted_data.append((text_content, slide_meta))

                # Presenter notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        notes_meta = {
                            "slide_number": slide_number,
                            "type": "presenter_notes",
                            "doc_type": "pptx",
                        }
                        formatted_notes = (
                            f"Presenter Notes (Slide {slide_number}):\n{notes_text}"
                        )
                        extracted_data.append((formatted_notes, notes_meta))

        except Exception as e:
            logger.error(f"[PPTX] Fatal error processing file {filepath}: {e}", exc_info=True)
            raise UnsupportedFileError(f"Error processing PPTX file {filepath}: {e}")

        return extracted_data
