from pptx import Presentation

slides = [
    "This is slide one.",
    "Slide two has more content.",
    "Third slide with important info.",
    "Final slide. Thanks!",
]

prs = Presentation()

for text in slides:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title & Content
    body = slide.shapes.placeholders[1].text_frame
    body.text = text

prs.save("tests/fixtures/demo.pptx")
print("Created tests/fixtures/demo.pptx with 4 slides.")
